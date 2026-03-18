import streamlit as st
import os
import re
import requests
import json
import time
import pandas as pd
import PyPDF2
import docx
import pptx
import urllib3
from io import BytesIO
from datetime import datetime
from fpdf import FPDF
from langchain_anthropic import ChatAnthropic
from langchain_google_genai import ChatGoogleGenerativeAI

# ==========================================
# 🚨 [추가된 코드] 사내 네트워크 SSL 인증서 에러 해결 🚨
# 이 코드는 다른 모듈들이 호출되기 전에 가장 먼저 실행되어야 합니다.
import ssl
import httpx

# 1. 파이썬 기본 SSL 검증 무력화 (권장하지 않지만 로컬 테스트용으로 가장 빠름)
ssl._create_default_https_context = ssl._create_unverified_context
# SSL 인증서 우회 시 발생하는 터미널 경고 메시지 숨기기
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# 2. httpx 통신 시 인증서 검증 무력화 환경변수 설정
os.environ['CURL_CA_BUNDLE'] = ''
os.environ['REQUESTS_CA_BUNDLE'] = ''
os.environ['SSL_CERT_FILE'] = ''
os.environ["HTTPX_SSL_VERIFY"] = "0"
# ==========================================

# --- 1. API 키 설정 ---
ANT_KEY = st.secrets.get("ANTHROPIC_API_KEY")
GEM_KEY = st.secrets.get("GOOGLE_API_KEY")
MANUS_API_KEY = st.secrets.get("MANUS_API_KEY")

# LLM 초기화 (모델명 고정)
fast_llm = ChatGoogleGenerativeAI(model="gemini-3-flash-preview", google_api_key=GEM_KEY)
smart_llm = ChatAnthropic(model="claude-sonnet-4-6", anthropic_api_key=ANT_KEY)

# --- 2. 초기 세션 상태 관리 ---
if "step" not in st.session_state:
    st.session_state.step = 1
for key in ["draft_content", "teaser_content", "design_recommendation", "month_title"]:
    if key not in st.session_state:
        st.session_state[key] = ""
if "manus_status" not in st.session_state:
    st.session_state.manus_status = "idle"
if "manus_url" not in st.session_state:
    st.session_state.manus_url = None

# --- 3. 유틸리티 함수 ---
def extract_text_from_file(uploaded_file):
    text = ""
    file_extension = uploaded_file.name.split(".")[-1].lower()
    try:
        if file_extension == "pdf":
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                extracted = page.extract_text()
                if extracted: text += extracted + "\n"
        elif file_extension in ["docx", "doc"]:
            doc = docx.Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_extension in ["pptx", "ppt"]:
            ppt = pptx.Presentation(uploaded_file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_extension in ["xlsx", "xls"]:
            df = pd.read_excel(uploaded_file)
            text += df.to_string() + "\n"
    except Exception as e:
        text += f"\n[{uploaded_file.name} 파일 읽기 오류: {e}]\n"
    return text

def create_manus_infographic(topic, report_content, style_instruction):
    if not MANUS_API_KEY: return None, "Manus API 키 누락"
    url = "https://api.manus.ai/v1/tasks"
    headers = {"API_KEY": MANUS_API_KEY, "Content-Type": "application/json"}
    
    enhanced_prompt = f"""
    시각적으로 돋보이는 16:9 비율의 전문적인 인포그래픽 프레젠테이션을 만들어주세요. 

    [🚨 필수 텍스트 표기 규칙]
    - 브랜드명 고정: 소식지 이름인 "expl'AI'n telink"는 반드시 소문자 바탕에 'AI'만 대문자로 표기해야 합니다.

    [디자인 가이드라인]
    ※ 가이드 라인은 절대로 수정하지 말고 그대로 반영
    1. 시각화 중심: 텍스트를 그대로 나열하기보다는, 내용을 한눈에 이해할 수 있도록 다이어그램, 차트, 인포그래픽 등을 적절히 활용해 주세요.
    2. 간결한 텍스트: 긴 문단보다는 핵심 키워드와 짧은 글머리 기호(Bullet points)를 사용하여 내용을 깔끔하게 요약해 주시면 좋습니다.
    3. 풍부한 시각 자료: 슬라이드 내용과 어울리는 고품질 이미지, 관련 아이콘 또는 일러스트를 자연스럽게 배치해 시각적인 매력을 더해 주세요.
    4. 세련된 레이아웃: 여백을 충분히 두고, 가독성이 좋은 폰트를 사용하여 매거진이나 테크 행사 발표 자료처럼 세련되고 깔끔한 느낌을 연출해 주세요.
    5. 아래 디자인 스타일대로만 슬라이드를 작성하고 절대 디자인 스타일을 벗어나지 말아주세요.

    주요 내용: {topic}
    디자인 스타일: {style_instruction}, 자연스러운 인포그래픽, 깔끔한 시각화
    본문 내용:
    {report_content[:4000]}
    """
    
    data = {
        "prompt": enhanced_prompt,
        "agentProfile": "manus-1.6-lite",
        "taskMode": "agent",
        "createShareableLink": True
    }
    
    try:
        response = requests.post(url, json=data, headers=headers, verify=False)
        if response.status_code not in [200, 201]: 
            return None, f"API 연결 실패 (상태 코드: {response.status_code})"
            
        res_json = response.json()
        task_id = res_json.get("task_id")
        task_url = res_json.get("task_url") # 🚨 추가: Manus 작업방 URL 확보
        
        if not task_id: return None, "Task ID를 발급받지 못했습니다."

        with st.status("📊 Manus 디자인 제작 중...", expanded=True):
            for _ in range(60):
                time.sleep(10)
                res = requests.get(f"{url}/{task_id}", headers={"API_KEY": MANUS_API_KEY}, verify=False).json()
                
                if res.get("status") == "completed":
                    files = res.get("files", [])
                    pptx_url = next((f.get('url') for f in files if isinstance(f, dict) and f.get('filename', '').endswith('.pptx')), None)
                    share_url = res.get("share_url")
                    
                    # 🚨 1순위: 파일 직접 다운로드, 2순위: 공유 링크, 3순위: 해당 작업의 Manus 대시보드
                    final_url = pptx_url or share_url or task_url
                    
                    if final_url: 
                        return final_url, "성공"
                    else: 
                        # 모든 링크가 없더라도 기본 Manus 홈페이지로 유도
                        return "https://manus.ai", "완료 (Manus 홈페이지에서 확인하세요)"
                
                elif res.get("status") in ["failed", "error"]:
                    return None, f"Manus API 내부 오류: {res.get('error', '알 수 없는 오류')}"
                    
        # 시간 초과 시에도 에러를 뱉지 않고 Manus 링크 제공
        return task_url or "https://manus.ai", "제작 시간이 길어지고 있습니다. Manus에서 확인하세요."
    except Exception as e: 
        return None, f"시스템 오류: {str(e)}"

def create_professional_pdf(text, title):
    pdf = FPDF()
    pdf.add_page()
    
    # 폰트 경로 설정 (현재 폴더의 나눔스퀘어 또는 윈도우 기본 맑은 고딕)
    eb_font = "NanumSquareEB.ttf"
    r_font = "NanumSquareR.ttf"
    windows_malgun = "C:/Windows/Fonts/malgun.ttf"
    
    # 제목용 폰트 로드
    if os.path.exists(eb_font):
        pdf.add_font('NS_EB', '', eb_font)
        t_f = 'NS_EB'
    elif os.path.exists(windows_malgun):
        pdf.add_font('Malgun', '', windows_malgun)
        t_f = 'Malgun'
    else:
        t_f = 'Helvetica'
        
    # 본문용 폰트 로드
    if os.path.exists(r_font):
        pdf.add_font('NS_R', '', r_font)
        b_f = 'NS_R'
    elif os.path.exists(windows_malgun):
        # 이미 로드된 맑은 고딕 재사용
        if t_f != 'Malgun':
            pdf.add_font('Malgun', '', windows_malgun)
        b_f = 'Malgun'
    else:
        b_f = 'Helvetica'
    
    # fpdf2는 UTF-8을 완벽 지원하므로 텍스트 인코딩 꼼수가 필요 없습니다.
    clean_title = title
    clean_text = text.replace('#', '').replace('*', '')
    
    # 제목 작성
    pdf.set_font(t_f, size=20)
    pdf.cell(0, 15, txt=clean_title, ln=1, align='L')
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(10)
    
    # 본문 작성
    pdf.set_font(b_f, size=11)
    pdf.multi_cell(0, 8, txt=clean_text)
    
    # PDF를 바이트 형태로 반환
    return bytes(pdf.output()), re.sub(r'[\\/*?:"<>|]', "", clean_title)

# --- 4. 에이전트 프롬프트 로직 ---
def extract_clean_text(content):
    if isinstance(content, list):
        return "".join([item.get("text", "") for item in content if isinstance(item, dict) and "text" in item])
    return str(content)

def generate_draft(data):
    inst_fin = "AI 보완 및 윤문 적용하여 내용 확장" if data['ai_fin'] else "입력된 텍스트 내용의 의미와 수치를 절대 변경하지 말고 그대로 슬라이드 텍스트로 배치할 것"
    inst_int = "AI 보완 및 윤문 적용하여 내용 확장" if data['ai_int'] else "입력된 텍스트 내용을 변경하지 말고 그대로 슬라이드 텍스트로 배치할 것"
    inst_insight = "AI 보완 및 윤문 적용하여 내용 확장" if data['ai_ins'] else "입력된 텍스트 내용을 변경하지 말고 그대로 배치할 것"
    inst_news = "AI 보완 및 윤문 적용하여 내용 확장" if data['ai_news_chk'] else "입력된 텍스트 내용을 변경하지 말고 그대로 배치할 것"

    prompt = f"""
    당신은 사내 소식지 "expl'AI'n telink: 텔링크를 말하다."의 전문 편집장입니다.
    이 출력물은 최종 보고서가 아니라, PPT 자동 생성기에 자동으로 전달될 '슬라이드 스크립트(Outline)'입니다.
    
    [🚨 필수 규칙 및 금지 사항]
    1. 브랜드명 고정: "expl'AI'n telink"는 반드시 소문자 바탕에 'AI'만 대문자로 표기하세요.
    2. 부가 설명 금지: 사용자에게 "이 내용을 복사해서 사용하세요"라거나 인사말 등의 불필요한 대화형 문구는 절대 출력하지 마세요. 오직 스크립트 내용만 출력해야 합니다.
    3. 형식: 화려한 마크다운 서식이나 표를 만들지 말고, 각 페이지별로 들어갈 '핵심 텍스트'만 간결하게 구분해서 작성하세요.

    [입력 정보 및 개별 지침]
    - 발행월: {data['month']}
    - 사내 실적: {data['financial']} (지침: {inst_fin})
    - 사내 주요 소식: {data['internal']} (지침: {inst_int})
    - AI Insight (방향성 및 첨부파일 문서 내용): {data['ai_insight']} (지침: {inst_insight})
    - AI 뉴스: {data['ai_news']} (지침: {inst_news})

    [출력 양식]
    페이지 1. 표지
    - 제목: expl'AI'n telink: 텔링크를 말하다. - {data['month']}
    - 안내 문구: ※ 이 문서는 AI 및 자동화 솔루션으로 제작 되었습니다.
    - 목차: 사내 소식 / AI Insight / 기타 소식

    페이지 2. 간지
    - 텍스트: 사내 소식 (재무성과 및 사내 주요 소식)

    페이지 3. 재무성과
    - (위 지침에 맞춰 실적 텍스트 작성)
    ※ 숫자를 강조하고 숫자 텍스트 크기를 다른 텍스트 크기보다 크게 표기, 필요 시 실적과 연관된 이미지 또는 아이콘 추가

    페이지 4. 사내 주요 소식
    - (위 지침에 맞춰 실적 텍스트 작성)
    ※ 해당 페이지는 각 본부별로 꾸밀 수 있는 내용과 연상되는 연관된 이미지 또는 아이콘 추가
    
    페이지 5. 간지
    - 텍스트: AI Insight

    페이지 6~X. AI Insight 본문
    - (위 지침에 맞춰 내용 작성, 필요시 여러 페이지로 분할)

    페이지 X+1. 간지
    - 텍스트: 기타 AI 소식

    페이지 X+2. AI 뉴스
    - (위 지침에 맞춰 내용 작성)

    페이지 X+3.
    - 티저 이미지 기획안 내용
    """
    
    res = smart_llm.invoke(prompt)
    return extract_clean_text(res.content)

def generate_teaser(ai_insight, ai_news, internal):
    prompt = f"""
    사내 소식지 "expl'AI'n telink"의 티저 이미지 기획안을 작성해 주세요. 
    
    [🚨 필수 조건 및 금지 사항]
    1. 표(Table) 작성 금지: 절대로 마크다운 표(|) 기호를 사용하지 말고, 깔끔한 줄글과 글머리 기호(•)로만 작성하세요. 텍스트를 너무 많이 담지 마세요.
    2. 레이아웃 비중 (매우 중요):
       - [메인 영역 - 70~80%]: 'AI Insight' 내용 중 가장 핵심적인 자극적인 카피 1~2개와 시각적 묘사 하나만 크게 제안하세요.
       - [서브 영역 - 20~30%]: 'AI 뉴스'와 '사내 주요 소식'은 아주 조그맣게 들어갑니다. 구체적인 내용은 다 빼고, 시선을 끄는 간단한 이모티콘 1개와 짧은 단어 수준의 핵심 텍스트(각 1~2줄)로만 극도로 요약하세요.
    3. 부가 설명 금지: 사용자에게 안내하는 문구나 인사말은 절대 넣지 마세요.

    [참고 데이터]
    - AI Insight: {ai_insight[:500]}
    - AI 뉴스: {ai_news[:300]}
    - 사내 주요 소식: {internal[:300]}

    [출력 양식 예시]
    ■ 메인 타이틀
    - 비주얼: [메인 배경 또는 그래픽 설명]
    - 메인 카피: "[강렬하고 짧은 한 줄 카피]"

    ■ 중앙 메인 영역 (AI Insight)
    - 비주얼 요소: [간단한 시각 요소 설명]
    - 카피: [1~2줄의 짧은 카피]

    ■ 하단 서브 영역 (사내 소식 & AI 뉴스)
    - 사내 소식: [이모티콘] [아주 짧은 핵심 단어/문구 1줄]
    - AI 뉴스: [이모티콘] [아주 짧은 핵심 단어/문구 1줄]
    """
    res = fast_llm.invoke(prompt)
    return extract_clean_text(res.content)

def revise_draft(current_draft, feedback):
    prompt = f"다음 초안을 사용자 피드백에 맞게 수정하되, PPT 스크립트(Outline) 형태의 간결한 구조는 반드시 유지해.\n\n[초안]\n{current_draft}\n\n[피드백]\n{feedback}"
    res = smart_llm.invoke(prompt)
    return extract_clean_text(res.content)

# 💡 [핵심] 사용자가 원하는 고퀄리티 프롬프트 구조 강제 적용
# 💡 [핵심] 촌스러운 디자인 배제 및 고퀄리티 프롬프트 강제 적용
def get_design_recommendation(month, ai_insight):
    prompt = f"""
    당신은 세계 최고의 프레젠테이션 아트 디렉터입니다.
    발행월({month})과 주요 다룰 내용(AI, 기술)을 바탕으로, **절대 뻔하지 않고 트렌디한** PPT 디자인 테마 3가지를 영문 프롬프트로 기획하세요.

    [🚨 엄격한 금지 사항 (Anti-Patterns)]
    - "Corporate", "Standard business", "Boring", "Basic blue tech", "Clean minimalist" 등 흔하고 지루한 비즈니스 템플릿 스타일은 절대 금지합니다.
    - 추상적이고 모호한 표현(예: "beautiful design", "modern look")을 쓰지 마세요.

    [✅ 필수 디자인 수준 및 포맷]
    - 사용자는 평범한 PPT가 아니라 매거진, 포스터, 감각적인 시각 자료를 원합니다. (예: Brutalism, Neo-Memphis, Cyberpunk, Editorial Zine, Claymorphism, Hand-drawn doodle 등)
    - 반드시 아래 <우수 예시>와 **동일한 구조(Design style, Layout, Color palette)와 디테일한 묘사**를 꽉꽉 채워야 합니다.

    <우수 예시>
    A hand-drawn doodle style infographic slide.

    Design style:
    - black marker sketch illustrations
    - yellow highlighter highlights
    - white crumpled paper background
    - casual whiteboard doodle drawing
    - simple stick figures and sketch icons
    - arrows, stars, speech bubbles, hand drawn shapes
    - handwritten typography style
    - playful startup presentation design

    Layout:
    - vertical infographic
    - two-column split comparison layout
    - section boxes with hand drawn borders
    - circled highlights and underlines
    - sketch diagrams and icons explaining ideas

    Color palette:
    black ink drawing + yellow highlight only
    minimalist visual storytelling infographic
    </우수 예시>

    주요 내용: {ai_insight[:300]}

    출력 형식:
    ### 1. [테마명 (예: Editorial Zine Style)]
    - **추천 이유**: [이유]
    - **프롬프트**:
    ```
    [여기에 예시 구조와 같이 작성된 상세 영문 프롬프트]
    ```
    - **참고**: [🎨 디자인 느낌 미리보기](https://www.google.com/search?tbm=isch&q=Editorial+Zine+presentation+design)

    ### 2. [테마명]
    (동일 구조)
    
    ### 3. [테마명]
    (동일 구조)
    """
    res = fast_llm.invoke(prompt)
    return extract_clean_text(res.content)

# --- 5. UI 영역 ---
st.set_page_config(page_title="expl'AI'n telink Studio", layout="wide")
st.title("📰 expl'AI'n telink 자동화 Agent")

# ---------------------------------------------------------
# Step 1: 정보 입력
# ---------------------------------------------------------
if st.session_state.step == 1:
    st.subheader("📝 Step 1. 이번 달 expl'AI'n telink 데이터 입력")
    
    now = datetime.now()
    default_month_str = f"{now.year}년 {now.month}월호"
    
    with st.form("input_form"):
        month = st.text_input("발행 월 (예: 2026년 3월호)", value=default_month_str)
        
        st.markdown("##### 📊 사내 실적 & AI Insight")
        row1_col1, row1_col2 = st.columns(2)
        with row1_col1:
            financial = st.text_area("1. 사내 실적 (매출/영업이익 등)", height=150)
            ai_fin = st.checkbox("이 항목에 AI 보완 적용하기", key="chk_fin", value=False)
            
        with row1_col2:
            ai_insight = st.text_area("3. AI Insight (핵심 요청사항이나 방향성을 입력)", height=100)
            uploaded_files = st.file_uploader("📂 AI Insight 참고 자료 첨부", type=["pdf", "docx", "pptx", "xlsx"], accept_multiple_files=True)
            ai_ins = st.checkbox("이 항목에 AI 보완 적용하기", key="chk_ins", value=True)

        st.markdown("---")
        st.markdown("##### 📰 사내 주요 소식 & AI 뉴스")
        row2_col1, row2_col2 = st.columns(2)
        with row2_col1:
            internal = st.text_area("2. 사내 주요 소식", height=150)
            ai_int = st.checkbox("이 항목에 AI 보완 적용하기", key="chk_int", value=False)

        with row2_col2:
            ai_news = st.text_area("4. AI 뉴스", height=150)
            ai_news_chk = st.checkbox("이 항목에 AI 보완 적용하기", key="chk_news", value=True)
            
        submit = st.form_submit_button("🚀 초안 및 디자인 추천 생성하기", use_container_width=True)
        
        if submit:
            if not month: st.warning("발행 월을 입력해주세요!")
            else:
                with st.spinner("첨부파일을 분석하고 초안과 맞춤형 PPT 디자인 테마를 기획 중입니다..."):
                    extracted_text = ""
                    if uploaded_files:
                        for file in uploaded_files:
                            extracted_text += f"\n--- [{file.name}] 내용 ---\n"
                            extracted_text += extract_text_from_file(file)
                    
                    combined_ai_insight = f"사용자 방향성: {ai_insight}\n\n[참고 문서 내용]\n{extracted_text}" if extracted_text else ai_insight
                    
                    data = {"month": month, "financial": financial, "ai_fin": ai_fin, 
                            "internal": internal, "ai_int": ai_int, "ai_insight": combined_ai_insight, 
                            "ai_ins": ai_ins, "ai_news": ai_news, "ai_news_chk": ai_news_chk}
                    
                    st.session_state.draft_content = generate_draft(data)
                    st.session_state.teaser_content = generate_teaser(combined_ai_insight, ai_news, internal)
                    st.session_state.design_recommendation = get_design_recommendation(month, combined_ai_insight)
                    st.session_state.month_title = month
                    st.session_state.step = 2
                    st.rerun()

# ---------------------------------------------------------
# Step 2: 초안 검토 및 수정 (직접 수정 기능 추가)
# ---------------------------------------------------------
elif st.session_state.step == 2:
    st.subheader("🧐 Step 2. Outline 초안 및 티저 직접 검토")
    st.info("아래 텍스트 박스에서 직접 내용을 수정하실 수 있습니다. 수정 후 우측 하단의 '저장 및 검토 완료'를 누르면 변경 사항이 PPT에 반영됩니다.")
    
    col_draft, col_teaser = st.columns([2, 1])
    with col_draft:
        st.markdown("### 📄 PPT 추출용 Outline 스크립트")
        edited_draft = st.text_area("초안 내용", value=st.session_state.draft_content, height=450, label_visibility="collapsed")
        
    with col_teaser:
        st.markdown("### 🎨 1페이지 티저 기획안")
        edited_teaser = st.text_area("티저 내용", value=st.session_state.teaser_content, height=450, label_visibility="collapsed")
    
    st.markdown("---")
    st.markdown("#### 🤖 AI에게 추가 수정 요청하기 (선택 사항)")
    feedback = st.chat_input("채팅으로 수정할 내용을 입력하세요 (예: 3페이지 실적 부분은 조금 더 요약해줘)")
    
    if feedback:
        with st.chat_message("user"): st.write(feedback)
        with st.spinner("피드백 반영 중..."):
            st.session_state.draft_content = revise_draft(edited_draft, feedback)
            st.session_state.teaser_content = edited_teaser
            st.rerun()
            
    st.divider()
    col_btn1, col_btn2 = st.columns(2)
    with col_btn1:
        if st.button("처음부터 다시 입력하기", use_container_width=True):
            st.session_state.step = 1; st.rerun()
    with col_btn2:
        if st.button("✅ 저장 및 검토 완료. 출력 단계로 이동", type="primary", use_container_width=True):
            st.session_state.draft_content = edited_draft
            st.session_state.teaser_content = edited_teaser
            st.session_state.step = 3; st.rerun()

# ---------------------------------------------------------
# Step 3: 최종 출력
# ---------------------------------------------------------
elif st.session_state.step == 3:
    st.subheader("🎉 Step 3. 최종 디자인 및 PPT 생성")
    
    # 🚨 [수정] 최종 텍스트를 세션에 저장하여 수정 가능하도록 만듭니다.
    if "final_full_text" not in st.session_state:
        st.session_state.final_full_text = f"{st.session_state.draft_content}\n\n[티저 기획안]\n{st.session_state.teaser_content}"
        
    report_title = f"expl'AI'n telink - {st.session_state.month_title}"
    
    st.success("✨ **AI가 제안하는 이번 달 맞춤형 디자인 테마 3선**")
    st.markdown(st.session_state.design_recommendation)
    
    st.divider()
    
    # --- 디자인 테마 원클릭 적용 및 직접 입력 버튼 ---
    if "selected_manus_style" not in st.session_state:
        st.session_state.selected_manus_style = ""

    prompts = re.findall(r'```(.*?)```', st.session_state.design_recommendation, re.DOTALL)
    
    if prompts:
        st.write("💡 **마음에 드는 테마를 클릭하거나, 직접 테마를 입력해 슬라이드 제작에 적용하세요!**")
        
        # 🚨 [수정] 추천 테마 개수 + 직접 입력 버튼 1개를 더해 컬럼 생성
        cols = st.columns(min(len(prompts) + 1, 4)) 
        
        for i, prompt_text in enumerate(prompts[:3]):
            with cols[i]:
                if st.button(f"🎨 추천 테마 {i+1} 적용", key=f"theme_btn_{i}", use_container_width=True):
                    st.session_state.selected_manus_style = prompt_text.strip()
                    st.rerun()
                    
        # 🚨 [추가] 4번째 버튼: 직접 입력 모드
        with cols[-1]:
            if st.button("✍️ 테마 직접 입력하기", type="primary", use_container_width=True):
                st.session_state.selected_manus_style = "" # 텍스트 박스 초기화
                st.rerun()
                
    st.divider()
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.write("📂 **1. 텍스트 내보내기 및 최종 편집**")
        st.info("아래 텍스트 박스에서 내용을 최종 수정한 후, 전체 복사(Ctrl+A, Ctrl+C)하여 붙여넣으세요. 수정한 내용은 PDF와 Manus 제작에 자동 반영됩니다.")
        
        # 🚨 [수정] st.code 대신 st.text_area를 사용하여 편집 가능하게 변경
        edited_final_text = st.text_area(
            "최종 스크립트 수정", 
            value=st.session_state.final_full_text, 
            height=450,
            label_visibility="collapsed"
        )
        
        # 사용자가 수정한 내용을 세션에 덮어쓰기
        st.session_state.final_full_text = edited_final_text
        
        # 다운로드되는 PDF도 수정된 텍스트 기반으로 굽기
        pdf_bytes, safe_name = create_professional_pdf(edited_final_text, report_title)
        st.download_button("📩 수정된 대본 PDF 다운로드", data=pdf_bytes, file_name=f"{safe_name}.pdf", use_container_width=True)
        
        if st.button("🔄 이전 단계로 돌아가기", use_container_width=True):
            st.session_state.step = 2; st.rerun()
            
    with col2:
        st.write("🎨 **시각화 자료**")
        
        if "manus_status" not in st.session_state:
            st.session_state.manus_status = "idle" 
        if "manus_url" not in st.session_state:
            st.session_state.manus_url = None

        with st.expander("📊 Manus 인포그래픽 슬라이드 제작", expanded=True):
            style_input = st.text_area(
                "디자인 테마 (버튼으로 선택하거나 직접 입력하세요)", 
                value=st.session_state.selected_manus_style,
                height=250,
                disabled=(st.session_state.manus_status == "processing")
            )

            if st.session_state.manus_status == "idle":
                if st.button("🚀 슬라이드 생성 시작", use_container_width=True):
                    st.session_state.selected_manus_style = style_input 
                    st.session_state.manus_status = "processing"
                    st.rerun()

            elif st.session_state.manus_status == "processing":
                st.button("⏳ 슬라이드 제작 중... 잠시만 기다려주세요", disabled=True, use_container_width=True)
                
                with st.status("📊 Manus 에이전트 가동 중...", expanded=True) as status:
                    # 🚨 Manus 제작 시에도 원본이 아닌 '수정된 텍스트(edited_final_text)'를 넘깁니다.
                    url, msg = create_manus_infographic(report_title, edited_final_text, style_input)
                    if url:
                        st.session_state.manus_url = url
                        st.session_state.manus_status = "completed"
                        status.update(label="✅ 제작 완료!", state="complete")
                        st.rerun()
                    else:
                        st.session_state.manus_status = "idle"
                        st.error(f"❌ 오류 발생: {msg}")
                        if st.button("🔄 다시 시도"): st.rerun()

            elif st.session_state.manus_status == "completed":
                st.link_button("📂 제작된 슬라이드 확인하기", st.session_state.manus_url, use_container_width=True, type="primary")
                if st.button("🆕 새로 만들기", use_container_width=True):
                    st.session_state.manus_status = "idle"
                    st.session_state.manus_url = None
                    st.rerun()
